#!/usr/bin/env python3
"""
Generate a comprehensive PowerPoint presentation for the VC UseCase Scoring Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    BLUE = RGBColor(31, 78, 121)
    LIGHT_BLUE = RGBColor(68, 114, 196)
    GREEN = RGBColor(112, 173, 71)
    ORANGE = RGBColor(237, 125, 49)
    GRAY = RGBColor(89, 89, 89)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "VC UseCase Scoring Platform"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Startup Idea Evaluation System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "October 2025"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = GRAY
    date_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Project Overview"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Purpose"
    p = tf.add_paragraph()
    p.text = "Connect startup aspirants with venture capitalists through AI-powered idea evaluation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Features"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "NLP-based idea-to-use-case matching with 3-dimensional scoring"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Role-based access (Startup Aspirants & VC Representatives)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Real-time evaluation with actionable recommendations"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Comprehensive reporting and analytics"
    p.level = 1
    
    # Slide 3: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Backend"
    p = tf.add_paragraph()
    p.text = "FastAPI 0.109 (Python 3.11)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "SQLAlchemy 2.0 ORM + PostgreSQL 16"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Redis 7 for caching"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "JWT authentication with bcrypt"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Frontend"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "React 18.2 + Vite 5.0"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Tailwind CSS 3.4"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Axios for API communication"
    p.level = 1
    
    # Slide 4: AI/ML Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AI/ML Technology Stack"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Natural Language Processing"
    p = tf.add_paragraph()
    p.text = "spaCy 3.7 - Entity recognition & linguistic features"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "NLTK 3.8 - Text preprocessing & tokenization"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sentence-Transformers 2.3 - Semantic embeddings"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Machine Learning"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Transformers 4.37 (Hugging Face)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "scikit-learn 1.4 - Cosine similarity & metrics"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "PyTorch 2.2 - Neural network backend"
    p.level = 1
    
    # Slide 5: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Microservices Architecture (Docker Compose)"
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Frontend Service (React)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Port 5173, Vite dev server"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Backend Service (FastAPI)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Port 8000, 15+ RESTful API endpoints"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "PostgreSQL Database"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Port 5433, 5 tables (User, VCUseCase, Idea, Evaluation, AuditLog)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Redis Cache"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Port 6380, NLP feature caching for performance"
    p.level = 1
    
    # Slide 6: Database Schema
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Database Schema"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Core Tables"
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "users - Authentication & role management"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "vc_usecases - VC investment criteria & pain points"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "ideas - Startup submissions with extracted features"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "evaluations - AI scoring results (alignment, novelty, viability)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "audit_logs - System activity tracking"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Relationships"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "User → Ideas (1:N)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "User → VCUseCases (1:N)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Idea + VCUseCase → Evaluation (N:M)"
    p.level = 1
    
    # Slide 7: API Endpoints
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "API Endpoints (15+)"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Authentication"
    p = tf.add_paragraph()
    p.text = "POST /api/register - User registration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "POST /api/token - JWT login"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "GET /api/users/me - Get current user"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Ideas (Startup Aspirants)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "POST /api/ideas - Submit startup idea"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "GET /api/ideas - List all ideas"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "VC Use Cases"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "POST /api/vc/usecases - Create use case"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "POST /api/vc/usecases/bulk - Bulk upload"
    p.level = 1
    
    # Slide 8: Scoring Algorithm
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AI Scoring Algorithm"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "3-Dimensional Scoring System"
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "1. Alignment Score (40%)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Semantic similarity between idea & use case"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Industry/domain matching"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Tag overlap analysis"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Novelty Score (30%)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Unique differentiators identification"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Innovation indicators"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Viability Score (30%)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Market signals analysis"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Problem-solution fit evaluation"
    p.level = 1
    
    # Slide 9: User Roles & Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "User Roles & Features"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Startup Aspirant Role"
    p = tf.add_paragraph()
    p.text = "Submit startup ideas with detailed descriptions"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "View AI-generated evaluations & scores"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Get actionable recommendations for improvement"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Track evaluation history"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "VC Representative Role"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Define investment criteria via use cases"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Bulk upload use cases (CSV/JSON)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Evaluate single or multiple ideas"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Generate comprehensive reports"
    p.level = 1
    
    # Slide 10: Deployment Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Deployment & DevOps"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Containerization"
    p = tf.add_paragraph()
    p.text = "Docker 28.5.1 with multi-stage builds"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Docker Compose for orchestration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Health checks for all services"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Volume mounting for development"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "CI/CD Pipeline"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "GitHub Actions workflow"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automated testing with pytest"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Code coverage reporting"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Linting & code quality checks"
    p.level = 1
    
    # Slide 11: Testing Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Testing & Quality Assurance"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Test Suite"
    p = tf.add_paragraph()
    p.text = "pytest 7.4.4 testing framework"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "pytest-asyncio for async endpoints"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "pytest-cov for coverage analysis"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Test Coverage"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "API endpoint testing (20+ tests)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Scoring engine unit tests"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Authentication & authorization tests"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Integration tests with test database"
    p.level = 1
    
    # Slide 12: Performance Optimization
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Performance Optimization"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Caching Strategy"
    p = tf.add_paragraph()
    p.text = "Redis caching for NLP feature extraction"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Reduces redundant ML computations"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Hash-based cache keys for ideas"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Database Optimization"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "SQLAlchemy query optimization"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Indexed columns for fast lookups"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Connection pooling"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Frontend Optimization"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Vite HMR for instant updates"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "React lazy loading & code splitting"
    p.level = 1
    
    # Slide 13: Security Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Security Implementation"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Authentication & Authorization"
    p = tf.add_paragraph()
    p.text = "JWT tokens with configurable expiration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "bcrypt password hashing (cost factor 12)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Role-based access control (RBAC)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "OAuth2 password flow"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "API Security"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "CORS configuration for frontend"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Input validation with Pydantic schemas"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "SQL injection protection via ORM"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Audit logging for security tracking"
    p.level = 1
    
    # Slide 14: Demo Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Demo: Test Results"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Test Scenario: AI Code Review Tool"
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Startup Idea Submitted"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Title: AI-Powered Code Review Tool"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Industry: Software Development"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Key Features: Real-time analysis, IDE integration, 95% accuracy"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "VC Use Case Matched"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Title: Developer Tools & Productivity"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Focus: AI-powered tools, automation, code quality"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "AI Evaluation Result: 40% Overall Score"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Indicates moderate alignment with VC criteria"
    p.level = 1
    
    # Slide 15: Project Statistics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Statistics"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Codebase Metrics"
    p = tf.add_paragraph()
    p.text = "60+ files created"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5,000+ lines of code"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "15+ API endpoints"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5 database models"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "7 frontend pages"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Dependencies"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "27 Python packages"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "15+ JavaScript libraries"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4 Docker containers"
    p.level = 1
    
    # Slide 16: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Enhancements"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Planned Features"
    p = tf.add_paragraph()
    p.text = "Advanced ML models (BERT, GPT integration)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Real-time collaboration between VCs and startups"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Market trend analysis dashboard"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Email notifications for evaluations"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Scalability Improvements"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Kubernetes deployment for production"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Horizontal scaling with load balancer"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Background job processing (Celery)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "CDN integration for frontend assets"
    p.level = 1
    
    # Slide 17: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Key Achievements"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(3.5))
    content_frame = content_box.text_frame
    
    p = content_frame.paragraphs[0]
    p.text = "✓ Full-stack application with modern tech stack"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "✓ AI/ML-powered scoring engine with NLP capabilities"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "✓ Dockerized microservices architecture"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "✓ Secure authentication & role-based access"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "✓ Comprehensive testing & CI/CD pipeline"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "✓ Production-ready deployment with Docker Compose"
    p.font.size = Pt(20)
    
    # Slide 18: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    contact_frame = contact_box.text_frame
    
    p = contact_frame.paragraphs[0]
    p.text = "Project URL: http://localhost:5173"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = GRAY
    
    p = contact_frame.add_paragraph()
    p.text = "API Docs: http://localhost:8000/docs"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = GRAY
    
    # Save presentation
    prs.save('VC_UseCase_Scoring_Platform_Presentation.pptx')
    print("✅ Presentation created successfully!")
    print("📄 File: VC_UseCase_Scoring_Platform_Presentation.pptx")
    print("📊 Total slides: 18")

if __name__ == "__main__":
    create_presentation()
